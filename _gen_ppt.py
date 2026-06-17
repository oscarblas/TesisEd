"""
Generador del PPT actualizado para Edwin
Tesis: Diseño del control de nivel para un sistema de cuatro tanques acoplados
mediante control predictivo
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============== ESTILO INSTITUCIONAL ==============
C_AZUL_PUCP   = RGBColor(0x00, 0x33, 0x66)
C_GRIS_OSCURO = RGBColor(0x33, 0x33, 0x33)
C_GRIS_CLARO  = RGBColor(0x88, 0x88, 0x88)
C_BLANCO      = RGBColor(0xFF, 0xFF, 0xFF)
C_ACENTO      = RGBColor(0xC0, 0x39, 0x2B)
C_VERDE       = RGBColor(0x2E, 0x7D, 0x32)

TITULO_TESIS_LINEA1 = "DISEÑO DEL CONTROL DE NIVEL PARA UN SISTEMA DE CUATRO"
TITULO_TESIS_LINEA2 = "TANQUES ACOPLADOS MEDIANTE CONTROL PREDICTIVO"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height

def add_blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_header(slide, capitulo_num=None, capitulo_titulo=None, seccion=None):
    """Añade encabezado superior y franja de capítulo."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_AZUL_PUCP
    bar.line.fill.background()
    txt = slide.shapes.add_textbox(Inches(0.3), Inches(0.05), Inches(13), Inches(0.45))
    tf = txt.text_frame
    tf.margin_left = 0; tf.margin_right = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "DISEÑO DEL CONTROL DE NIVEL PARA UN SISTEMA DE CUATRO TANQUES ACOPLADOS MEDIANTE CONTROL PREDICTIVO"
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = C_BLANCO

    if capitulo_num is not None:
        cap = slide.shapes.add_textbox(Inches(0.3), Inches(0.7), Inches(12.5), Inches(0.4))
        cap_tf = cap.text_frame
        cap_tf.margin_left = 0
        p_cap = cap_tf.paragraphs[0]
        r_cap = p_cap.add_run()
        r_cap.text = f"Capítulo {capitulo_num}.  {capitulo_titulo}"
        r_cap.font.size = Pt(15); r_cap.font.bold = True; r_cap.font.color.rgb = C_AZUL_PUCP
    if seccion:
        sec = slide.shapes.add_textbox(Inches(0.3), Inches(1.1), Inches(12.5), Inches(0.35))
        p_sec = sec.text_frame.paragraphs[0]
        r_sec = p_sec.add_run()
        r_sec.text = seccion
        r_sec.font.size = Pt(13); r_sec.font.italic = True; r_sec.font.color.rgb = C_GRIS_OSCURO

def add_footer(slide, numero):
    foot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, H - Inches(0.3), W, Inches(0.3))
    foot.fill.solid()
    foot.fill.fore_color.rgb = C_AZUL_PUCP
    foot.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(12.5), H - Inches(0.3), Inches(0.7), Inches(0.3))
    p = tx.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = str(numero)
    r.font.size = Pt(10); r.font.color.rgb = C_BLANCO; r.font.bold = True

def add_bullets(slide, items, left, top, width, height, size=14, color=C_GRIS_OSCURO,
                spacing=0.8):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = "•  " + it
        r.font.size = Pt(size)
        r.font.color.rgb = color
        p.space_after = Pt(int(spacing * 8))
    return tx

def add_image_placeholder(slide, left, top, width, height, label):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    box.line.color.rgb = C_GRIS_CLARO; box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "(IMAGEN)"
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = C_ACENTO
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label
    r2.font.size = Pt(10); r2.font.color.rgb = C_GRIS_OSCURO; r2.font.italic = True

def add_titulo_seccion(slide, texto, top, size=18, color=None):
    if color is None: color = C_AZUL_PUCP
    tx = slide.shapes.add_textbox(Inches(0.3), top, Inches(12.5), Inches(0.5))
    p = tx.text_frame.paragraphs[0]
    r = p.add_run(); r.text = texto
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = color

def add_box_destacado(slide, left, top, width, height, titulo, items, color_fondo, color_texto=None):
    if color_texto is None: color_texto = C_GRIS_OSCURO
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = color_fondo
    box.line.color.rgb = color_fondo
    tf = box.text_frame
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1); tf.margin_bottom = Inches(0.1)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = titulo
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = color_texto
    for it in items:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = "  ◦  " + it
        r2.font.size = Pt(11); r2.font.color.rgb = color_texto

# =====================================================================
# SLIDE 1 — PORTADA
# =====================================================================
s = add_blank_slide()
# Fondo
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
bg.fill.solid(); bg.fill.fore_color.rgb = C_AZUL_PUCP; bg.line.fill.background()

# Banda blanca central
banda = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), W, Inches(2.7))
banda.fill.solid(); banda.fill.fore_color.rgb = C_BLANCO; banda.line.fill.background()

# Universidad arriba
tx = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.6))
p = tx.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "PONTIFICIA UNIVERSIDAD CATÓLICA DEL PERÚ"
r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = C_BLANCO

tx2 = s.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.5))
p = tx2.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "FACULTAD DE CIENCIAS E INGENIERÍA"
r.font.size = Pt(14); r.font.color.rgb = C_BLANCO

# Título central
tx = s.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(12.3), Inches(0.7))
p = tx.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = TITULO_TESIS_LINEA1
r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = C_AZUL_PUCP
tx = s.shapes.add_textbox(Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.7))
p = tx.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = TITULO_TESIS_LINEA2
r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = C_AZUL_PUCP

# Trabajo de tesis 2
tx = s.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.5))
p = tx.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Trabajo de Tesis 2  —  Semestre 2026-1"
r.font.size = Pt(14); r.font.italic = True; r.font.color.rgb = C_GRIS_OSCURO

# Autor y asesor
tx = s.shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.45))
p = tx.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Autor:  Edwin Jesús Davila Avalos"
r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = C_BLANCO

tx = s.shapes.add_textbox(Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.45))
p = tx.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Asesor:  Juan Javier Sotomayor Moriano"
r.font.size = Pt(13); r.font.color.rgb = C_BLANCO

tx = s.shapes.add_textbox(Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.4))
p = tx.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Lima, junio 2026"
r.font.size = Pt(12); r.font.color.rgb = C_BLANCO

# =====================================================================
# SLIDE 2 — ÍNDICE
# =====================================================================
s = add_blank_slide()
add_header(s)
add_footer(s, 2)

add_titulo_seccion(s, "ÍNDICE", Inches(0.85), size=22)

# Columna izquierda
items_left = [
    ("Capítulo 1.", "Estudio del control de nivel en sistemas de tanques acoplados"),
    ("Capítulo 2.", "Modelado del sistema y fundamentos del Control Predictivo GPC"),
    ("Capítulo 3.", "Diseño del controlador predictivo GPC MIMO"),
]
items_right = [
    ("Capítulo 4.", "Análisis comparativo del GPC frente al PI con desacoplador"),
    ("", "Conclusiones generales"),
    ("", "Trabajo futuro"),
    ("", "Bibliografía"),
]

yy = Inches(1.7)
for cap, txt in items_left:
    tx = s.shapes.add_textbox(Inches(0.7), yy, Inches(6.0), Inches(0.7))
    p = tx.text_frame.paragraphs[0]
    if cap:
        r = p.add_run(); r.text = cap + "  "
        r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = C_AZUL_PUCP
    r2 = p.add_run(); r2.text = txt
    r2.font.size = Pt(14); r2.font.color.rgb = C_GRIS_OSCURO
    yy += Inches(0.85)

yy = Inches(1.7)
for cap, txt in items_right:
    tx = s.shapes.add_textbox(Inches(7.0), yy, Inches(6.0), Inches(0.7))
    p = tx.text_frame.paragraphs[0]
    if cap:
        r = p.add_run(); r.text = cap + "  "
        r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = C_AZUL_PUCP
    r2 = p.add_run(); r2.text = txt
    r2.font.size = Pt(14); r2.font.color.rgb = C_GRIS_OSCURO
    yy += Inches(0.85)

# =====================================================================
# SLIDE 3 — CAP 1: Importancia y problemática
# =====================================================================
s = add_blank_slide()
add_header(s, 1, "Estudio del control de nivel en sistemas de tanques acoplados",
           "1.1 Importancia y problemática")
add_footer(s, 3)

add_bullets(s, [
    "Control de nivel: una de las operaciones más recurrentes en la industria",
    "Sectores: petroquímica, minería, alimentaria, farmacéutica, energética",
    "Sistemas MIMO con acoplamiento ⇒ control SISO clásico es insuficiente",
    "Limitaciones: no linealidades, restricciones físicas, perturbaciones",
], Inches(0.4), Inches(1.6), Inches(6.3), Inches(4), size=14)

add_image_placeholder(s, Inches(7.0), Inches(1.6), Inches(5.9), Inches(2.6),
                      "Ejemplo industrial — control de pH en circuito de cal")
add_image_placeholder(s, Inches(7.0), Inches(4.3), Inches(5.9), Inches(2.7),
                      "Diagrama de sistema TITO con acoplamiento")

# =====================================================================
# SLIDE 4 — CAP 1: Sistema de 4 tanques (Johansson + Lab PUCP)
# =====================================================================
s = add_blank_slide()
add_header(s, 1, "Estudio del control de nivel en sistemas de tanques acoplados",
           "1.2 Sistema de cuatro tanques acoplados")
add_footer(s, 4)

add_image_placeholder(s, Inches(0.4), Inches(1.6), Inches(4.0), Inches(3.0),
                      "Esquema Johansson (1998)")
add_image_placeholder(s, Inches(0.4), Inches(4.7), Inches(4.0), Inches(2.2),
                      "Planta piloto Lab. Control Avanzado PUCP")

# Box: fase mínima vs no mínima
add_titulo_seccion(s, "Fase del sistema según γ₁ + γ₂", Inches(1.6), size=14)
add_box_destacado(s, Inches(4.7), Inches(2.1), Inches(8.3), Inches(1.0),
                  "Fase MÍNIMA (P⁻):  1 < γ₁+γ₂ < 2",
                  ["Ceros en semiplano izquierdo  |  Comportamiento estable",
                   "Caminos directos dominan  |  Control predictivo aplicable"],
                  RGBColor(0xE8, 0xF5, 0xE9))
add_box_destacado(s, Inches(4.7), Inches(3.25), Inches(8.3), Inches(1.0),
                  "Fase NO MÍNIMA (P⁺):  0 < γ₁+γ₂ < 1",
                  ["Cero en semiplano derecho  |  Respuesta inversa transitoria",
                   "Limita el ancho de banda alcanzable por cualquier controlador lineal"],
                  RGBColor(0xFF, 0xEB, 0xEE))

add_image_placeholder(s, Inches(4.7), Inches(4.4), Inches(8.3), Inches(2.6),
                      "Comparación de la respuesta al escalón en fase mínima vs no mínima")

# =====================================================================
# SLIDE 5 — CAP 1: Estado del arte + Objetivos
# =====================================================================
s = add_blank_slide()
add_header(s, 1, "Estudio del control de nivel en sistemas de tanques acoplados",
           "1.3 Estado del arte  |  1.5 Objetivos")
add_footer(s, 5)

# Estado del arte
add_titulo_seccion(s, "Técnicas aplicadas en la literatura", Inches(1.6), size=14)
add_bullets(s, [
    "Control clásico:  PID descentralizado, PI con desacoplador",
    "Control avanzado:  fuzzy-PID, control adaptativo, robusto fraccional",
    "Inteligencia artificial:  PI-NN, PID-NN con PSO",
    "Control predictivo:  DMC, GPC, NMPC, DMPC",
], Inches(0.4), Inches(2.0), Inches(6.0), Inches(2.5), size=12)

# Objetivos
add_titulo_seccion(s, "Objetivo general", Inches(1.6), size=14)
add_box_destacado(s, Inches(7.0), Inches(2.0), Inches(6.0), Inches(1.2),
                  "Diseñar un controlador GPC multivariable",
                  ["Regular niveles de los tanques inferiores",
                   "Considerar acoplamientos, restricciones y perturbaciones",
                   "Comparar con control PI clásico con desacoplador"],
                  RGBColor(0xE3, 0xF2, 0xFD))

add_titulo_seccion(s, "Objetivos específicos", Inches(3.4), size=14)
add_bullets(s, [
    "Caracterizar la dinámica MIMO del sistema",
    "Obtener modelo linealizado en el punto de operación",
    "Diseñar el algoritmo GPC con restricciones",
    "Implementar en MATLAB/Simulink",
    "Comparar GPC vs PI+Desacoplador",
], Inches(0.4), Inches(3.8), Inches(12.5), Inches(3.5), size=12, spacing=0.5)

# =====================================================================
# SLIDE 6 — CAP 2: Modelado de la planta
# =====================================================================
s = add_blank_slide()
add_header(s, 2, "Modelado del sistema y fundamentos del Control Predictivo GPC",
           "2.1 Modelamiento del sistema de cuatro tanques acoplados")
add_footer(s, 6)

add_bullets(s, [
    "Modelo no lineal:  balance de masa + ley de Bernoulli",
    "4 ecuaciones diferenciales acopladas con términos √h",
    "Linealización en torno al punto de operación h₀ = (h₁⁰, h₂⁰, 25, 25)",
    "Representación en espacio de estados  +  matriz de transferencia G(s)",
], Inches(0.4), Inches(1.7), Inches(5.8), Inches(4), size=13)

add_image_placeholder(s, Inches(6.4), Inches(1.7), Inches(6.6), Inches(2.7),
                      "Esquema del sistema con γ₁, γ₂ y caminos directo/cruzado")

add_image_placeholder(s, Inches(6.4), Inches(4.5), Inches(6.6), Inches(2.5),
                      "Ecuaciones no lineales del modelo (sistema de 4 EDOs)")

# =====================================================================
# SLIDE 7 — CAP 2: GPC y modelo CARIMA
# =====================================================================
s = add_blank_slide()
add_header(s, 2, "Modelado del sistema y fundamentos del Control Predictivo GPC",
           "2.2 Control Predictivo Generalizado (GPC)")
add_footer(s, 7)

add_titulo_seccion(s, "Tres pilares del GPC", Inches(1.5), size=14)

add_box_destacado(s, Inches(0.4), Inches(2.0), Inches(4.1), Inches(1.6),
                  "1. Modelo de predicción",
                  ["Modelo CARIMA",
                   "Polinomios A, B, C en z⁻¹",
                   "Operador Δ = 1 − z⁻¹  →  acción integral"],
                  RGBColor(0xE3, 0xF2, 0xFD))

add_box_destacado(s, Inches(4.6), Inches(2.0), Inches(4.1), Inches(1.6),
                  "2. Función de costo",
                  ["Cuadrática:  error² + esfuerzo²",
                   "Pesos δ (seguimiento) y λ (control)",
                   "Horizontes N₁, N₂, Nu"],
                  RGBColor(0xFF, 0xF8, 0xE1))

add_box_destacado(s, Inches(8.8), Inches(2.0), Inches(4.1), Inches(1.6),
                  "3. Ley de control",
                  ["Solución óptima via QP",
                   "Δu* = (GᵀQG + R)⁻¹ Gᵀ Q (W − F)",
                   "Solo se aplica el primer Δu"],
                  RGBColor(0xE8, 0xF5, 0xE9))

add_image_placeholder(s, Inches(2.5), Inches(3.9), Inches(8.5), Inches(3.0),
                      "Estructura básica del MPC: predicción, optimización y horizonte deslizante")

# =====================================================================
# SLIDE 8 — CAP 2: GPC MIMO + validación del modelo
# =====================================================================
s = add_blank_slide()
add_header(s, 2, "Modelado del sistema y fundamentos del Control Predictivo GPC",
           "2.3 GPC MIMO  |  2.4 Validación del modelo")
add_footer(s, 8)

add_titulo_seccion(s, "GPC en sistemas multivariables", Inches(1.5), size=14)
add_bullets(s, [
    "Extensión a polinomios matriciales A(z⁻¹), B(z⁻¹)",
    "Matriz dinámica G por bloques  G_ij  (uno por par entrada-salida)",
    "Función de costo con matrices Q y R bloque-diagonales",
    "Misma estructura analítica que el caso SISO",
], Inches(0.4), Inches(1.95), Inches(6.0), Inches(2.8), size=12, spacing=0.4)

add_titulo_seccion(s, "Validación del modelo: lineal vs no lineal", Inches(4.5), size=14)
add_bullets(s, [
    "Comparación de respuesta al escalón sobre ambos modelos",
    "Métrica FIT% para cuantificar la similitud",
    "Validez confirmada en una vecindad del punto de operación",
], Inches(0.4), Inches(5.0), Inches(6.0), Inches(2.0), size=12, spacing=0.4)

add_image_placeholder(s, Inches(6.7), Inches(1.95), Inches(6.3), Inches(2.5),
                      "Esquema GPC MIMO: matriz G por bloques")
add_image_placeholder(s, Inches(6.7), Inches(4.6), Inches(6.3), Inches(2.4),
                      "Comparación lineal vs no lineal en operación normal — FIT > 95%")

# =====================================================================
# SLIDE 9 — CAP 3: Selección de Ts y N
# =====================================================================
s = add_blank_slide()
add_header(s, 3, "Diseño del controlador GPC para el sistema TITO",
           "3.3.2 Selección del horizonte y tiempo de muestreo")
add_footer(s, 9)

add_titulo_seccion(s, "Parámetros del MODELO (fijados a priori)", Inches(1.6), size=14)
add_bullets(s, [
    "Constantes de tiempo dominantes:  τ_dom ≈ 30 s",
    "Criterio:  T_s = min(0.1 · τ_ij)  →  sobre-muestreo por carácter MIMO",
    "Criterio:  N · T_s ≥ 5 · τ_dom  para cubrir la dinámica completa",
    "Estos parámetros NO se sintonizan: dependen del proceso",
], Inches(0.4), Inches(2.05), Inches(7.5), Inches(2.5), size=13, spacing=0.4)

# Caja con parámetros adoptados
add_box_destacado(s, Inches(8.2), Inches(1.95), Inches(4.8), Inches(2.5),
                  "Valores adoptados (Tabla 3.1)",
                  ["T_s = 1 s",
                   "N   = 50",
                   "N · T_s = 50 s (cubre > 1.5 τ_dom)",
                   "Coherente con Åström & Wittenmark",
                   "para sistemas MIMO ruidosos"],
                  RGBColor(0xE3, 0xF2, 0xFD))

add_image_placeholder(s, Inches(0.4), Inches(4.7), Inches(12.5), Inches(2.4),
                      "Coeficientes de respuesta al escalón g_ij[k] para los 4 subprocesos")

# =====================================================================
# SLIDE 10 — CAP 3: Matriz G y respuesta libre F
# =====================================================================
s = add_blank_slide()
add_header(s, 3, "Diseño del controlador GPC para el sistema TITO",
           "3.3.3 Matriz dinámica G  |  3.3.4 Respuesta libre F")
add_footer(s, 10)

add_box_destacado(s, Inches(0.4), Inches(1.7), Inches(6.2), Inches(2.8),
                  "Matriz dinámica G",
                  ["Por bloques  G_ij  (uno por par entrada-salida)",
                   "Cada G_ij es triangular inferior",
                   "Coeficientes:  respuesta al escalón g_ij[k]",
                   "Estructura: G = [G₁₁ G₁₂; G₂₁ G₂₂]"],
                  RGBColor(0xE3, 0xF2, 0xFD))

add_box_destacado(s, Inches(6.7), Inches(1.7), Inches(6.3), Inches(2.8),
                  "Respuesta libre F",
                  ["Trayectoria con Δu = 0  (modelo libre)",
                   "Propagación del modelo discreto con u(k−1) fijo",
                   "Vector F apila las predicciones futuras",
                   "Predicción total:  Ŷ = G · ΔU + F"],
                  RGBColor(0xFF, 0xF8, 0xE1))

add_image_placeholder(s, Inches(0.4), Inches(4.7), Inches(12.5), Inches(2.4),
                      "Estructura de la matriz G por bloques + diagrama de la predicción")

# =====================================================================
# SLIDE 11 — CAP 3: Restricciones (QP)
# =====================================================================
s = add_blank_slide()
add_header(s, 3, "Diseño del controlador GPC para el sistema TITO",
           "3.3.7 Tratamiento de restricciones (QP)")
add_footer(s, 11)

add_titulo_seccion(s, "Tres tipos de restricciones físicas", Inches(1.6), size=14)
add_bullets(s, [
    "Sobre los incrementos:  Δu_min ≤ Δu(k+i) ≤ Δu_max",
    "Sobre los valores absolutos:  u_min ≤ u(k+i) ≤ u_max",
    "Sobre las salidas (opcional):  y_min ≤ y(k+i) ≤ y_max",
], Inches(0.4), Inches(2.0), Inches(7.0), Inches(2.0), size=13, spacing=0.5)

add_titulo_seccion(s, "Forma estándar QP", Inches(4.0), size=14)
add_box_destacado(s, Inches(0.4), Inches(4.4), Inches(8.0), Inches(1.2),
                  "min  ½ ΔUᵀ H ΔU + fᵀ ΔU     s.a.    A_ineq · ΔU ≤ b_ineq",
                  ["H = 2 (GᵀQG + R)",
                   "f = 2 GᵀQ (F − W)",
                   "Resolución con quadprog (Algoritmo: active-set)"],
                  RGBColor(0xE3, 0xF2, 0xFD))

add_image_placeholder(s, Inches(8.6), Inches(2.0), Inches(4.4), Inches(5.0),
                      "Diagrama de flujo del controlador GPC")

# =====================================================================
# SLIDE 12 — CAP 3: Métodos de sintonización
# =====================================================================
s = add_blank_slide()
add_header(s, 3, "Diseño del controlador GPC para el sistema TITO",
           "3.4 Sintonización del controlador GPC MIMO")
add_footer(s, 12)

add_titulo_seccion(s, "Cuatro métodos comparados  (uno por familia)", Inches(1.6), size=14)

add_box_destacado(s, Inches(0.4), Inches(2.1), Inches(6.2), Inches(1.4),
                  "1. Clarke-Mohtadi (1987)  —  Heurístico-analítico",
                  ["Reglas prácticas de los autores originales del GPC",
                   "Nu pequeño (1-3), δ unitario, λ moderado"],
                  RGBColor(0xE3, 0xF2, 0xFD))

add_box_destacado(s, Inches(6.7), Inches(2.1), Inches(6.2), Inches(1.4),
                  "2. Shridhar-Cooper (1997)  —  Analítico explícito",
                  ["Aproximación FOPDT por par entrada-salida",
                   "Fórmulas cerradas para Nu y λ"],
                  RGBColor(0xFF, 0xF8, 0xE1))

add_box_destacado(s, Inches(0.4), Inches(3.7), Inches(6.2), Inches(1.4),
                  "3. PSO  —  Metaheurístico global",
                  ["Optimización por enjambre de partículas",
                   "Exploración global sobre [log₁₀(λ), Nu]"],
                  RGBColor(0xE8, 0xF5, 0xE9))

add_box_destacado(s, Inches(6.7), Inches(3.7), Inches(6.2), Inches(1.4),
                  "4. Nelder-Mead (fminsearch)  —  Numérico directo",
                  ["Búsqueda local sin gradiente",
                   "Refina el óptimo entregado por PSO"],
                  RGBColor(0xFC, 0xE4, 0xEC))

add_titulo_seccion(s, "Comparación mediante score combinado", Inches(5.4), size=14)
add_bullets(s, [
    "6 métricas:  IAE, ISE, ITAE, t_est, sobrepico, esfuerzo",
    "Normalización min-max y suma ponderada  →  Score (menor = mejor)",
    "Ts = 1 s  y  N = 50  permanecen fijos durante toda la sintonización",
], Inches(0.4), Inches(5.85), Inches(12.5), Inches(1.5), size=12, spacing=0.3)

# =====================================================================
# SLIDE 13 — CAP 3: Resultado de la sintonización
# =====================================================================
s = add_blank_slide()
add_header(s, 3, "Diseño del controlador GPC para el sistema TITO",
           "3.4.6 Selección del método ganador")
add_footer(s, 13)

add_image_placeholder(s, Inches(0.4), Inches(1.6), Inches(7.2), Inches(5.4),
                      "Gráfica de barras del score combinado por método")

add_titulo_seccion(s, "Parámetros finales del GPC", Inches(1.6), size=14)
add_box_destacado(s, Inches(7.8), Inches(2.0), Inches(5.2), Inches(2.4),
                  "Sintonización ganadora",
                  ["T_s     =  1 s",
                   "N       =  50",
                   "N_u     =  5",
                   "δ       =  [10, 10]",
                   "λ       =  0.0077  (valor aproximado)"],
                  RGBColor(0xE8, 0xF5, 0xE9))

add_titulo_seccion(s, "Observaciones", Inches(4.6), size=14)
add_bullets(s, [
    "El método de optimización numérica ofrece el mejor compromiso",
    "Nu = 5 reduce el sobreajuste sin sacrificar agilidad",
    "Parámetros aplicados al controlador final del Cap. 4",
], Inches(7.8), Inches(5.0), Inches(5.2), Inches(2), size=11, spacing=0.4)

# =====================================================================
# SLIDE 14 — CAP 4: Estrategia comparativa
# =====================================================================
s = add_blank_slide()
add_header(s, 4, "Análisis comparativo del GPC frente al PI con desacoplador",
           "4.1 Estrategia comparativa")
add_footer(s, 14)

add_titulo_seccion(s, "Pregunta central del capítulo", Inches(1.6), size=14)
add_box_destacado(s, Inches(0.4), Inches(2.0), Inches(12.5), Inches(0.9),
                  "¿En qué condiciones se justifica el GPC frente al PI clásico con desacoplador?",
                  [],
                  RGBColor(0xE3, 0xF2, 0xFD))

add_titulo_seccion(s, "Controladores comparados", Inches(3.2), size=14)

add_box_destacado(s, Inches(0.4), Inches(3.7), Inches(6.2), Inches(3.0),
                  "GPC multivariable (Cap. 3)",
                  ["Modelo CARIMA matricial",
                   "Resolución de QP cada T_s",
                   "Predicción a N pasos con horizonte deslizante",
                   "Restricciones explícitas en u y Δu",
                   "Coordinación nativa de los dos lazos"],
                  RGBColor(0xE3, 0xF2, 0xFD))

add_box_destacado(s, Inches(6.7), Inches(3.7), Inches(6.3), Inches(3.0),
                  "PI descentralizado + Desacoplador",
                  ["Pairing:  u₁ → h₄,  u₂ → h₃",
                   "Sintonización IMC:  K_p = τ / (K · λ_imc)",
                   "Anti-windup por saturación incremental",
                   "Desacoplador estático Skogestad:  D = [1 −k₁₂; −k₂₁ 1]",
                   "Cancela acople SOLO en estado estacionario"],
                  RGBColor(0xFF, 0xF8, 0xE1))

# =====================================================================
# SLIDE 15 — CAP 4: Escenario integrado
# =====================================================================
s = add_blank_slide()
add_header(s, 4, "Análisis comparativo del GPC frente al PI con desacoplador",
           "4.3 Escenario integrado de simulación")
add_footer(s, 15)

add_titulo_seccion(s, "Trayectoria de eventos sobre la planta NO LINEAL", Inches(1.5), size=14)

add_box_destacado(s, Inches(0.4), Inches(2.0), Inches(6.2), Inches(0.85),
                  "t = 0  →  400 s  :  Arranque desde h = 0", [],
                  RGBColor(0xE3, 0xF2, 0xFD))
add_box_destacado(s, Inches(0.4), Inches(2.95), Inches(6.2), Inches(0.85),
                  "t = 400 s  :  Cambio SP h₃   25 → 30  cm", [],
                  RGBColor(0xE8, 0xF5, 0xE9))
add_box_destacado(s, Inches(0.4), Inches(3.90), Inches(6.2), Inches(0.85),
                  "t = 800 s  :  Cambio SP h₄   25 → 20  cm", [],
                  RGBColor(0xE8, 0xF5, 0xE9))
add_box_destacado(s, Inches(0.4), Inches(4.85), Inches(6.2), Inches(0.85),
                  "t = 1100 s  :  Inyección de ruido gaussiano (σ = 0.3 cm)", [],
                  RGBColor(0xFF, 0xF8, 0xE1))
add_box_destacado(s, Inches(0.4), Inches(5.80), Inches(6.2), Inches(0.85),
                  "t = 1200 s  :  h₃ → 25  y  h₄ → 35  (lejos del punto)", [],
                  RGBColor(0xFC, 0xE4, 0xEC))

add_image_placeholder(s, Inches(6.9), Inches(2.0), Inches(6.1), Inches(5.0),
                      "Trayectoria de referencias h₃(t) y h₄(t)")

# =====================================================================
# SLIDE 16 — CAP 4: Resultados ante acoplamiento
# =====================================================================
s = add_blank_slide()
add_header(s, 4, "Análisis comparativo del GPC frente al PI con desacoplador",
           "4.4 Resultados ante acoplamiento cruzado")
add_footer(s, 16)

add_image_placeholder(s, Inches(0.4), Inches(1.6), Inches(8.4), Inches(5.5),
                      "Curvas comparativas h₃(t) y h₄(t) — GPC vs PI+Desacoplador")

add_box_destacado(s, Inches(9.0), Inches(1.7), Inches(4.0), Inches(2.5),
                  "GPC",
                  ["Anticipa el acople",
                   "Coordina ambas bombas",
                   "Interacción cruzada mínima"],
                  RGBColor(0xE3, 0xF2, 0xFD))
add_box_destacado(s, Inches(9.0), Inches(4.4), Inches(4.0), Inches(2.5),
                  "PI + Desacoplador",
                  ["Desacopla solo en estacionario",
                   "Perturbación visible en el otro lazo",
                   "INT_ij claramente mayor"],
                  RGBColor(0xFF, 0xF8, 0xE1))

# =====================================================================
# SLIDE 17 — CAP 4: Resultado clave (h4 = 35)
# =====================================================================
s = add_blank_slide()
add_header(s, 4, "Análisis comparativo del GPC frente al PI con desacoplador",
           "4.5 Resultado clave  —  operación lejos del punto")
add_footer(s, 17)

# Banner destacado
add_box_destacado(s, Inches(0.4), Inches(1.6), Inches(12.5), Inches(1.0),
                  "RESULTADO PRINCIPAL DE LA TESIS",
                  ["En t ≥ 1200 s, el PI con desacoplador NO alcanza h₄ = 35 cm; el GPC SÍ lo logra."],
                  C_AZUL_PUCP, color_texto=C_BLANCO)

add_image_placeholder(s, Inches(0.4), Inches(2.85), Inches(8.0), Inches(4.2),
                      "Detalle del tramo t ≥ 1200 s — h₄ vs referencia 35 cm")

add_box_destacado(s, Inches(8.6), Inches(2.85), Inches(4.4), Inches(2.0),
                  "¿Por qué falla el PI?",
                  ["IMC sintonizado en h₀ = 25",
                   "τ depende de √h  →  cambia con h",
                   "Desacoplador estático válido SOLO en h₀"],
                  RGBColor(0xFF, 0xF8, 0xE1))
add_box_destacado(s, Inches(8.6), Inches(5.0), Inches(4.4), Inches(2.0),
                  "¿Por qué SÍ llega el GPC?",
                  ["Horizonte de predicción a N pasos",
                   "Acción integral del CARIMA",
                   "Coordinación coordinada de las bombas"],
                  RGBColor(0xE8, 0xF5, 0xE9))

# =====================================================================
# SLIDE 18 — CAP 4: Métricas comparativas
# =====================================================================
s = add_blank_slide()
add_header(s, 4, "Análisis comparativo del GPC frente al PI con desacoplador",
           "4.4 Métricas comparativas")
add_footer(s, 18)

add_titulo_seccion(s, "Dos ventanas de evaluación", Inches(1.6), size=14)
add_bullets(s, [
    "Métricas GLOBALES  →  toda la simulación (incluye arranque)",
    "Métricas en OPERACIÓN NORMAL  →  a partir de t = 400 s",
    "Métrica adicional INT_ij  →  cuantifica acoplamiento cruzado",
], Inches(0.4), Inches(2.0), Inches(12.0), Inches(1.5), size=13, spacing=0.4)

add_image_placeholder(s, Inches(0.4), Inches(3.6), Inches(6.2), Inches(3.4),
                      "Tabla comparativa de métricas (IAE, ISE, ITAE, t_est, Mp, esfuerzo)")

add_image_placeholder(s, Inches(6.7), Inches(3.6), Inches(6.3), Inches(3.4),
                      "Gráfico de barras: GPC vs PI+Desacoplador para cada métrica")

# =====================================================================
# SLIDE 19 — CAP 4: Implementación en Simulink
# =====================================================================
s = add_blank_slide()
add_header(s, 4, "Análisis comparativo del GPC frente al PI con desacoplador",
           "4.6 Implementación en Simulink")
add_footer(s, 19)

add_titulo_seccion(s, "Verificación adicional sobre Simulink", Inches(1.5), size=14)
add_bullets(s, [
    "Ambos controladores implementados como MATLAB Functions",
    "Planta no lineal en un único bloque con realimentación por Unit Delay",
    "Solver:  ode45 variable-step  (para preservar precisión MIMO)",
    "Resultados consistentes con los scripts de MATLAB",
], Inches(0.4), Inches(1.95), Inches(6.6), Inches(2.6), size=12, spacing=0.4)

add_image_placeholder(s, Inches(0.4), Inches(4.5), Inches(6.6), Inches(2.6),
                      "Diagrama Simulink del PI + Desacoplador")

add_image_placeholder(s, Inches(7.2), Inches(1.95), Inches(5.8), Inches(5.2),
                      "Diagrama Simulink del GPC")

# =====================================================================
# SLIDE 20 — CONCLUSIONES
# =====================================================================
s = add_blank_slide()
add_header(s)
add_footer(s, 20)

add_titulo_seccion(s, "CONCLUSIONES", Inches(0.85), size=22)

add_bullets(s, [
    "Se diseñó un controlador GPC multivariable para el sistema TITO de cuatro tanques bajo formulación CARIMA + ecuaciones diofánticas.",
    "Se compararon 4 métodos de sintonización (Clarke-Mohtadi, Shridhar-Cooper, PSO, Nelder-Mead) mediante un score combinado de seis métricas.",
    "El GPC reduce significativamente el acoplamiento cruzado entre lazos respecto al PI con desacoplador (métrica INT_ij).",
    "El PI con desacoplador falla en regiones alejadas del punto de operación (h₄ = 35 cm); el GPC sí responde gracias a su predicción a N pasos.",
    "Ambos controladores fueron validados en MATLAB y replicados en Simulink, demostrando reproducibilidad.",
    "Trade-off identificado: mayor desempeño del GPC vs mayor costo computacional y dependencia del modelo.",
], Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.5), size=13, spacing=0.7)

# =====================================================================
# SLIDE 21 — TRABAJO FUTURO
# =====================================================================
s = add_blank_slide()
add_header(s)
add_footer(s, 21)

add_titulo_seccion(s, "TRABAJO FUTURO", Inches(0.85), size=22)

add_bullets(s, [
    "Implementación experimental en la planta piloto del Lab. de Control Avanzado de la PUCP.",
    "Evaluación del controlador frente a perturbaciones físicas reales (fugas, fallos de actuadores, ruido industrial).",
    "Estudio de variantes:  GPC adaptativo, NMPC para manejo explícito de la no linealidad.",
    "Análisis de robustez ante incertidumbre paramétrica y envejecimiento del proceso.",
    "Integración con plataformas SCADA / MES industriales para validación en entorno productivo.",
], Inches(0.5), Inches(1.7), Inches(12.3), Inches(5), size=13, spacing=0.8)

# =====================================================================
# SLIDE 22 — BIBLIOGRAFÍA (parte 1)
# =====================================================================
s = add_blank_slide()
add_header(s)
add_footer(s, 22)
add_titulo_seccion(s, "BIBLIOGRAFÍA  (1/2)", Inches(0.85), size=22)

bibs1 = [
    "[1]  E. Laubwald, “Coupled tank system,” Control Systems Principles, 2005.",
    "[2]  W. Yuan, “Mathematical Model Analysis and Control Strategy,” ICMCR 2023.",
    "[3]  Short, M. & Selvakumar, A. (2020). Non-Linear Tank Level Control for Industrial Applications.",
    "[4]  Numsomran, V., Tipsuwanporn, V., Tirasesth, K. “Modeling of the modified Quadruple-Tank Process,” SICE 2008.",
    "[5]  Azam, S.N.M. & Jørgensen, J.B. “Modeling and simulation of a modified quadruple tank system,” ICCSCE 2015.",
    "[6]  Yu, Y. et al. (2024). Método de control cooperativo para sistemas acoplados. Scientific Reports.",
    "[7]  Albertos, P. & Sala, A. Multivariable Control Systems. Springer, 2004.",
    "[8]  K. H. Johansson, “The quadruple-tank process,” IEEE Trans. Control Syst. Tech., 2000.",
    "[9]  Pugliese, L. et al. (2022). Planta didáctica de bajo coste para sistemas multivariables.",
    "[10] Sánchez Zurita, V.A. (2018). Diseño de un sistema de control predictivo multivariable. Tesis MSc PUCP.",
    "[11] Gouta, H. et al. (2022). Anti-disturbance composite tracking control. Automatika.",
    "[12] Camacho, E.F. & Bordons, C. Model Predictive Control, 2nd ed., Springer, 2007.",
]
add_bullets(s, bibs1, Inches(0.5), Inches(1.6), Inches(12.5), Inches(5.5), size=11, spacing=0.2)

# =====================================================================
# SLIDE 23 — BIBLIOGRAFÍA (parte 2)
# =====================================================================
s = add_blank_slide()
add_header(s)
add_footer(s, 23)
add_titulo_seccion(s, "BIBLIOGRAFÍA  (2/2)", Inches(0.85), size=22)

bibs2 = [
    "[13] Clarke, D.W. “Application of Generalized Predictive Control to Industrial Processes,” IEEE Control Syst. Mag., 1988.",
    "[14] Santana, H.G. et al. “Application of Multivariable PID Controllers in a Coupled Tank System,” INDUSCON 2018.",
    "[15] Choudhary, P.K., Raj, P. & Das, D.K. “Controller Design for Decoupled TITO Coupled Tank System,” ICSPCRE 2024.",
    "[16] Albin Raj, R.J. & Deepa, S.N. “Modeling and implementation of various controllers used for Quadruple-Tank,” ICCPCT 2016.",
    "[17] Rivera, D.E., Morari, M. & Skogestad, S. (1986). “Internal Model Control: PID Controller Design.”",
    "[18] Åström, K.J. & Hägglund, T. (2006). “Advanced PID Control.” ISA.",
    "[19] Skogestad, S. & Postlethwaite, I. (2005). Multivariable Feedback Control. Wiley.",
    "[20] Shridhar, R. & Cooper, D.J. (1997). “A Tuning Strategy for Unconstrained Multivariable Model Predictive Control.”",
    "[21] Eberhart, R. & Kennedy, J. (1995). “A new optimizer using particle swarm theory,” MHS.",
    "[22] Åström, K.J. & Wittenmark, B. (1997). Computer-Controlled Systems. Prentice Hall.",
    "[23] Ogata, K. (2010). Modern Control Engineering. Prentice Hall.",
]
add_bullets(s, bibs2, Inches(0.5), Inches(1.6), Inches(12.5), Inches(5.5), size=11, spacing=0.2)

# =====================================================================
# SLIDE 24 — GRACIAS
# =====================================================================
s = add_blank_slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
bg.fill.solid(); bg.fill.fore_color.rgb = C_AZUL_PUCP; bg.line.fill.background()

tx = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5))
p = tx.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "MUCHAS GRACIAS"
r.font.size = Pt(54); r.font.bold = True; r.font.color.rgb = C_BLANCO

tx = s.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.8))
p = tx.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "¿Preguntas?"
r.font.size = Pt(24); r.font.italic = True; r.font.color.rgb = C_BLANCO

tx = s.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.5))
p = tx.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Edwin Jesús Davila Avalos  —  Asesor: Juan Javier Sotomayor Moriano  —  Lima, 2026"
r.font.size = Pt(13); r.font.color.rgb = C_BLANCO

prs.save("C:\\Users\\HP\\OneDrive\\Escritorio\\edwin\\EdwinDavila_Presentacion_tesis_actualizada.pptx")
print(f"PPT generada con {len(prs.slides)} diapositivas")
